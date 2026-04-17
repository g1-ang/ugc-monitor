"""
build_phase3_pptx.py
────────────────────
phase3_results.json (매치된 UGC 유저) + phase3_candidates.json (이미지 URL)를 합쳐
PPTX 리포트로 생성.

실행:
  python scripts/build_phase3_pptx.py
  python scripts/build_phase3_pptx.py --out custom_report.pptx
"""

from __future__ import annotations
import os, io, json, argparse, requests
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 팔레트 (main.py와 동일) ────────────────
TEAL   = RGBColor(0x5B, 0xBF, 0xAD); TEAL_D = RGBColor(0x3D, 0x9E, 0x8E)
AMBER  = RGBColor(0xE0, 0x9A, 0x5A); AMBER_D= RGBColor(0xC0, 0x78, 0x40)
BLUE   = RGBColor(0x6A, 0x9F, 0xD8); BLUE_D = RGBColor(0x48, 0x78, 0xB8)
DARK   = RGBColor(0x1C, 0x1C, 0x1A); MED    = RGBColor(0x4A, 0x4A, 0x48)
GRAY   = RGBColor(0x8A, 0x88, 0x80); LGRAY  = RGBColor(0xC4, 0xC2, 0xBA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF); BG     = RGBColor(0xF7, 0xF6, 0xF2)
SURF   = RGBColor(0xFF, 0xFF, 0xFF); BORDER = RGBColor(0xE2, 0xE0, 0xD8)

TYPE_CFG = {
    "feed":    {"ko": "피드",   "c": TEAL,  "d": TEAL_D,  "bg": RGBColor(0xE4,0xF5,0xF2)},
    "story":   {"ko": "스토리", "c": AMBER, "d": AMBER_D, "bg": RGBColor(0xFD,0xF0,0xE2)},
    "profile": {"ko": "프사",   "c": BLUE,  "d": BLUE_D,  "bg": RGBColor(0xE6,0xF0,0xFB)},
}


def txt(slide, text, l, t, w, h, size, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.name = "Arial"
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color


def rct(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line
    else:    s.line.fill.background()
    return s


def enrich_results(confirmed_ugc: list, candidates: list) -> list:
    """confirmed_ugc에 image_url 추가 (candidates에서 lookup)"""
    cand_by_user = {u["username"]: u for u in candidates}
    out = []
    for r in confirmed_ugc:
        uname = r["username"]
        utype = r["ugc_type"]
        link  = r.get("feed_url", "") or ""
        cand  = cand_by_user.get(uname, {})

        image_url = ""
        if utype == "profile":
            image_url = cand.get("profile_url", "")
        elif utype == "story":
            stories = cand.get("story_image_urls") or []
            if not stories and cand.get("story_image_url"):
                stories = [cand["story_image_url"]]
            if stories:
                image_url = stories[0]
        elif utype == "feed":
            # post_url로 매칭되는 feed item 찾기
            for item in cand.get("latest_feed_items", []) or []:
                if item.get("post_url") == link:
                    image_url = item.get("image_url", "")
                    break
            if not image_url:
                feed_items = cand.get("latest_feed_items") or []
                if feed_items and feed_items[0].get("image_url"):
                    image_url = feed_items[0]["image_url"]
                elif cand.get("latest_feed_urls"):
                    image_url = cand["latest_feed_urls"][0]

        out.append({
            "username":    uname,
            "type":        utype,
            "link":        link,
            "image_url":   image_url,
            "detected_at": "",
        })
    return out


def build_pptx(results: list, stats: dict, date_str: str, total: int) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1: Cover ────────────────────────────
    s1 = prs.slides.add_slide(blank)
    s1.background.fill.solid(); s1.background.fill.fore_color.rgb = BG
    rct(s1, 0, 0, prs.slide_width, Inches(0.08), TEAL)
    rct(s1, Inches(1.0), Inches(0.5), Inches(0.06), Inches(1.5), TEAL)
    txt(s1, "UGC MONITORING REPORT — Phase 3",
        Inches(1.2), Inches(0.55), Inches(10), Inches(0.4), size=10, color=GRAY)
    txt(s1, "pitapat_prompt",
        Inches(1.2), Inches(0.95), Inches(10), Inches(0.9), size=36, bold=True, color=DARK)
    txt(s1, f"{date_str}  ·  대상 {total}명",
        Inches(1.2), Inches(1.9), Inches(10), Inches(0.4), size=13, color=GRAY)
    rct(s1, Inches(1.0), Inches(2.85), Inches(11.3), Inches(0.015), BORDER)

    card_data = [
        ("피드 UGC",   stats.get("feed",    0), "feed"),
        ("스토리 UGC", stats.get("story",   0), "story"),
        ("프사 변경",  stats.get("profile", 0), "profile"),
    ]
    for (label, count, ttype), left in zip(card_data, [Inches(1.0), Inches(4.6), Inches(8.2)]):
        cfg = TYPE_CFG[ttype]
        rct(s1, left, Inches(3.1), Inches(3.3), Inches(2.5), SURF, BORDER)
        rct(s1, left, Inches(3.1), Inches(3.3), Inches(0.07), cfg["c"])
        txt(s1, label, left+Inches(0.22), Inches(3.32), Inches(2.9), Inches(0.4),
            size=11, color=GRAY)
        txt(s1, str(count), left+Inches(0.18), Inches(3.72), Inches(2.9), Inches(1.2),
            size=60, bold=True, color=cfg["d"])
        txt(s1, "건", left+Inches(0.22), Inches(4.95), Inches(2.9), Inches(0.4),
            size=13, color=GRAY)
    txt(s1, f"총  {len(results)}건  감지",
        Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.55),
        size=15, bold=True, color=MED, align=PP_ALIGN.CENTER)

    # ── Slides 2+: Per-UGC ────────────────────────
    for r in results:
        sl = prs.slides.add_slide(blank)
        sl.background.fill.solid(); sl.background.fill.fore_color.rgb = BG
        cfg   = TYPE_CFG.get(r.get("type", "feed"), TYPE_CFG["feed"])
        uname = r.get("username", "")
        link  = r.get("link") or ""

        rct(sl, 0, 0, prs.slide_width, Inches(0.07), cfg["c"])

        # 이미지
        img_ok = False
        img_url = r.get("image_url", "")
        if img_url:
            try:
                ir = requests.get(img_url, timeout=10,
                                  headers={"User-Agent": "Mozilla/5.0"})
                if ir.status_code == 200:
                    rct(sl, Inches(0.56), Inches(0.41), Inches(7.4), Inches(6.9), LGRAY)
                    sl.shapes.add_picture(io.BytesIO(ir.content),
                                          Inches(0.5), Inches(0.35),
                                          Inches(7.4), Inches(6.9))
                    img_ok = True
            except Exception:
                pass
        if not img_ok:
            rct(sl, Inches(0.5), Inches(0.35), Inches(7.4), Inches(6.9), BORDER)
            txt(sl, "이미지 없음 또는 로드 실패",
                Inches(0.5), Inches(3.5), Inches(7.4), Inches(0.5),
                size=14, color=LGRAY, align=PP_ALIGN.CENTER)

        PNL = Inches(8.15); PW = Inches(4.7)
        rct(sl, PNL, Inches(0.45), Inches(1.3), Inches(0.42), cfg["bg"], cfg["c"])
        txt(sl, cfg["ko"], PNL, Inches(0.46), Inches(1.3), Inches(0.40),
            size=12, bold=True, color=cfg["d"], align=PP_ALIGN.CENTER)
        txt(sl, f"@{uname}", PNL, Inches(1.1), PW, Inches(0.75),
            size=26, bold=True, color=DARK)
        txt(sl, f"https://instagram.com/{uname}",
            PNL, Inches(1.95), PW, Inches(0.4), size=11, color=GRAY)
        rct(sl, PNL, Inches(2.55), PW, Inches(0.018), BORDER)
        if link:
            txt(sl, "게시물 링크", PNL, Inches(2.75), PW, Inches(0.35),
                size=10, color=LGRAY)
            txt(sl, link, PNL, Inches(3.1), PW, Inches(0.55),
                size=11, color=cfg["d"])
        else:
            txt(sl, "링크 없음 (스토리/프사)", PNL, Inches(2.75), PW, Inches(0.4),
                size=11, color=LGRAY, italic=True)
        rct(sl, PNL, Inches(6.9), PW, Inches(0.018), BORDER)
        txt(sl, {"feed":"피드 게시물","story":"스토리","profile":"프로필 사진"}.get(
                r.get("type","feed"),""),
            PNL, Inches(7.0), PW, Inches(0.35), size=10, color=LGRAY)

    buf = io.BytesIO()
    prs.save(buf); buf.seek(0)
    return buf


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--results",    default="phase3_results.json")
    parser.add_argument("--candidates", default="phase3_candidates.json")
    parser.add_argument("--out",        default=None)
    args = parser.parse_args()

    with open(args.results, encoding="utf-8") as f:
        res = json.load(f)
    with open(args.candidates, encoding="utf-8") as f:
        cands = json.load(f)

    confirmed = res.get("confirmed_ugc", [])
    enriched  = enrich_results(confirmed, cands)

    stats = {"feed": 0, "story": 0, "profile": 0}
    for r in enriched:
        t = r["type"]
        if t in stats: stats[t] += 1

    date_str = datetime.now().strftime("%Y-%m-%d")
    out = args.out or f"phase3_report_{datetime.now().strftime('%Y%m%d')}.pptx"

    print(f"매치된 유저: {len(confirmed)}명")
    print(f"  · feed:    {stats['feed']}")
    print(f"  · story:   {stats['story']}")
    print(f"  · profile: {stats['profile']}")
    print(f"이미지 다운로드 + PPTX 생성 중...")

    buf = build_pptx(enriched, stats, date_str, total=len(res.get("all_results", [])))
    with open(out, "wb") as f:
        f.write(buf.getvalue())
    print(f"✅ 저장: {out} ({os.path.getsize(out)//1024} KB)")


if __name__ == "__main__":
    main()
