"""
build_presentation_pptx.py
──────────────────────────
바이브코딩 동아리 발표용 PPTX 생성.
사용자가 calamus.design에서 디자인 재작업할 예정이므로 텍스트 위주로 단순하게.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK   = RGBColor(0x1C, 0x1C, 0x1A)
MED    = RGBColor(0x4A, 0x4A, 0x48)
GRAY   = RGBColor(0x8A, 0x88, 0x80)
ACCENT = RGBColor(0x3D, 0x9E, 0x8E)  # teal
RED    = RGBColor(0xC0, 0x44, 0x44)
BG     = RGBColor(0xF7, 0xF6, 0xF2)


def txt(slide, text, l, t, w, h, size, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    paras = text.split("\n")
    for i, line in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = "Pretendard"  # 한글 깔끔
        run.font.size = Pt(size); run.font.bold = bold
        run.font.italic = italic; run.font.color.rgb = color


def bullets(slide, items, l, t, w, h, size=14, color=MED, gap=1.3):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"·  {item}"
        run.font.name = "Pretendard"
        run.font.size = Pt(size); run.font.color.rgb = color
        p.space_after = Pt(size * 0.5)


def make():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1: Title ──────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    txt(s, "VIBE CODING · 2026",
        Inches(1), Inches(2.2), Inches(11), Inches(0.4), 12, color=GRAY)
    txt(s, "Claude와 함께한 바이브코딩 실전기",
        Inches(1), Inches(2.6), Inches(11), Inches(1.2), 40, bold=True, color=DARK)
    txt(s, "AI UGC 모니터링 자동화 파이프라인 구축",
        Inches(1), Inches(3.8), Inches(11), Inches(0.7), 22, color=ACCENT)
    txt(s, "수동 모니터링의 한계를 넘어, Vision AI로 찾는 우리 브랜드의 흔적",
        Inches(1), Inches(4.6), Inches(11), Inches(0.5), 14, color=MED, italic=True)
    txt(s, "pitapat_prompt",
        Inches(1), Inches(6.5), Inches(11), Inches(0.4), 12, color=GRAY)

    # ── Slide 2: Problem ────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    txt(s, "01 · Problem",
        Inches(1), Inches(0.8), Inches(11), Inches(0.4), 12, color=ACCENT, bold=True)
    txt(s, "왜 만들었나?",
        Inches(1), Inches(1.3), Inches(11), Inches(1.0), 36, bold=True, color=DARK)
    txt(s, "수동 모니터링의 한계",
        Inches(1), Inches(2.4), Inches(11), Inches(0.5), 16, color=MED)

    bullets(s, [
        "수백 명의 댓글 작성자 프로필을 일일이 확인 → 불가능 (노동 집약적)",
        "우리 AI 프롬프트로 생성된 이미지를 올린 유저(UGC)를 놓치고 있음",
        "이 UGC 유저를 발굴해서 브랜드 인게이지먼트 확인 필요",
        "\"사람 눈 대신 AI 눈으로 자동화하자\"",
    ], Inches(1), Inches(3.2), Inches(11), Inches(3.5), size=17)

    # ── Slide 3: Architecture ───────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    txt(s, "02 · Architecture",
        Inches(1), Inches(0.8), Inches(11), Inches(0.4), 12, color=ACCENT, bold=True)
    txt(s, "4단계 자동화 파이프라인",
        Inches(1), Inches(1.3), Inches(11), Inches(1.0), 32, bold=True, color=DARK)

    steps = [
        ("STEP 1", "수집", "인스타그램 게시물 댓글 데이터 추출 (xlsx)"),
        ("STEP 2", "스캔", "Apify로 댓글러의 프사 · 피드 · 스토리 자동 수집"),
        ("STEP 3", "판별", "Vision AI로 레퍼런스 이미지와 비교 (우리 프롬프트로 만들었는가?)"),
        ("STEP 4", "시각화", "대시보드 + Google Sheets에 결과 기록"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        top = Inches(2.6 + i * 0.95)
        txt(s, num, Inches(1), top, Inches(1.2), Inches(0.4), 13, bold=True, color=ACCENT)
        txt(s, title, Inches(2.3), top, Inches(2.0), Inches(0.4), 17, bold=True, color=DARK)
        txt(s, desc, Inches(4.5), top, Inches(8), Inches(0.4), 14, color=MED)

    # ── Slide 4: Tech Stack ─────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    txt(s, "03 · Tech Stack",
        Inches(1), Inches(0.8), Inches(11), Inches(0.4), 12, color=ACCENT, bold=True)
    txt(s, "프로젝트를 지탱하는 기술들",
        Inches(1), Inches(1.3), Inches(11), Inches(1.0), 32, bold=True, color=DARK)

    stack = [
        ("FRONTEND",  "웹 대시보드 (HTML · JavaScript)",
         "레퍼런스 이미지 업로드 · 프롬프트 입력 · 실시간 진행률 표시"),
        ("BACKEND",   "Python · FastAPI",
         "스캔 파이프라인 오케스트레이션 · Background Task 처리"),
        ("AI (Vision)", "Gemini 2.0 Flash (via NAMC Vertex AI)",
         "이미지-이미지 비교 및 텍스트 프롬프트 결합 판별"),
        ("SCRAPING",  "Apify Actor (Instagram Profile/Story/Post Scraper)",
         "댓글러의 프사 · 스토리 · 피드 수집"),
        ("DATA",      "Google Sheets API",
         "유저 목록 · 스캔 기록 · 결과 저장"),
    ]
    for i, (label, tech, desc) in enumerate(stack):
        top = Inches(2.4 + i * 0.82)
        txt(s, label, Inches(1), top, Inches(1.8), Inches(0.4), 11, bold=True, color=ACCENT)
        txt(s, tech,  Inches(2.9), top, Inches(10), Inches(0.4), 14, bold=True, color=DARK)
        txt(s, desc,  Inches(2.9), top + Inches(0.38), Inches(10), Inches(0.4), 12, color=GRAY)

    # ── Slide 5: Troubleshooting TOP 5 ──────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    txt(s, "04 · Troubleshooting",
        Inches(1), Inches(0.8), Inches(11), Inches(0.4), 12, color=ACCENT, bold=True)
    txt(s, "삽질 TOP 5",
        Inches(1), Inches(1.3), Inches(11), Inches(1.0), 32, bold=True, color=DARK)
    txt(s, "\"버그는 항상 가장 단순한 곳에 있다\"",
        Inches(1), Inches(2.25), Inches(11), Inches(0.4), 13, color=GRAY, italic=True)

    issues = [
        ("1", "Google Sheets 시트명 누락",
         "!ugc_users 접두어 빠뜨려서 엉뚱한 탭에 기록됨 — 찾는 데 2시간"),
        ("2", "이미지 URL vs 포스트 URL 혼동",
         "AI에게 post URL을 전달 → 이미지 못 불러서 판별 실패"),
        ("3", "모델 세 번 교체",
         "Gemini → CLIP → Qwen → 최종 Gemini 2.0 Flash (API 한도 · 속도 · 정확도 순차 해결)"),
        ("4", "정확도 11%에서 출발",
         "56개 false positive · 하이브리드로 87.5%까지 올림"),
        ("5", "NAMC는 NAVER 사내망 전용",
         "Render 배포가 애초에 불가능했음 — 로컬 도구로 전환"),
    ]
    for i, (n, title, desc) in enumerate(issues):
        top = Inches(2.85 + i * 0.78)
        txt(s, n,     Inches(1), top, Inches(0.6), Inches(0.5), 20, bold=True, color=RED)
        txt(s, title, Inches(1.7), top, Inches(10), Inches(0.4), 15, bold=True, color=DARK)
        txt(s, desc,  Inches(1.7), top + Inches(0.38), Inches(10), Inches(0.4), 12, color=MED)

    # ── Slide 6: AI Logic (Hybrid) ─────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    txt(s, "05 · AI Logic",
        Inches(1), Inches(0.8), Inches(11), Inches(0.4), 12, color=ACCENT, bold=True)
    txt(s, "하이브리드 판별 — 이 프로젝트의 핵심 깨달음",
        Inches(1), Inches(1.3), Inches(11), Inches(1.0), 28, bold=True, color=DARK)

    txt(s, "BEFORE", Inches(1), Inches(2.7), Inches(4), Inches(0.4), 13, bold=True, color=RED)
    txt(s, "레퍼런스 이미지 3장 비교", Inches(1), Inches(3.1), Inches(4.5), Inches(0.4), 15, bold=True, color=DARK)
    bullets(s, [
        "이미지만으로 \"비슷한지\" 판단",
        "3장 중 1장이라도 매치 → YES",
        "\"AI 생성 스타일이 비슷함\"에 과도하게 반응",
        "→ 정확도 11% (63개 매치 중 7개 정답)",
    ], Inches(1), Inches(3.7), Inches(5.5), Inches(3), size=12)

    txt(s, "AFTER", Inches(7.2), Inches(2.7), Inches(4), Inches(0.4), 13, bold=True, color=ACCENT)
    txt(s, "이미지 + 프롬프트 텍스트 + 다수결", Inches(7.2), Inches(3.1), Inches(5.5), Inches(0.4), 15, bold=True, color=DARK)
    bullets(s, [
        "레퍼런스 이미지 + 원본 AI 프롬프트 텍스트 함께 전달",
        "Vision AI가 \"이 프롬프트로 만들었나?\"를 직접 판별",
        "3장 중 2장 이상 매치해야 YES (다수결 투표)",
        "→ 정확도 87.5% · false positive 56개 → 1개",
    ], Inches(7.2), Inches(3.7), Inches(5.5), Inches(3), size=12, color=DARK)

    # ── Slide 7: Efficiency ─────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    txt(s, "06 · Efficiency",
        Inches(1), Inches(0.8), Inches(11), Inches(0.4), 12, color=ACCENT, bold=True)
    txt(s, "비용과 속도의 균형",
        Inches(1), Inches(1.3), Inches(11), Inches(1.0), 32, bold=True, color=DARK)

    bullets(s, [
        "Apify 스토리 배치 5 → 20 : 호출 횟수 1/4로 절감",
        "Rate Limit 회피 : 동시 처리 5 → 2 + 글로벌 세마포어 (6개 제한)",
        "실패 사용자만 재시도 : CDN 만료 URL은 직접 다운로드 후 전송",
        "처음 전체 스캔 처리 시간 : 913명 × 약 30분 → 10분",
        "1회 스캔 비용 : 약 $2-3 (주로 Apify 스크래핑 비용)",
    ], Inches(1), Inches(2.7), Inches(11.3), Inches(4.5), size=16)

    # ── Slide 8: Result ─────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    txt(s, "07 · Result",
        Inches(1), Inches(0.8), Inches(11), Inches(0.4), 12, color=ACCENT, bold=True)
    txt(s, "실제 결과",
        Inches(1), Inches(1.3), Inches(11), Inches(1.0), 32, bold=True, color=DARK)

    # Big numbers
    txt(s, "913", Inches(1), Inches(2.7), Inches(4), Inches(1.5), 80, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(s, "분석 대상", Inches(1), Inches(4.2), Inches(4), Inches(0.4), 13, color=GRAY, align=PP_ALIGN.CENTER)

    txt(s, "7", Inches(4.8), Inches(2.7), Inches(3.8), Inches(1.5), 80, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(s, "실제 UGC 매치", Inches(4.8), Inches(4.2), Inches(3.8), Inches(0.4), 13, color=GRAY, align=PP_ALIGN.CENTER)

    txt(s, "87.5%", Inches(8.5), Inches(2.7), Inches(4), Inches(1.5), 80, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(s, "정확도 (Precision)", Inches(8.5), Inches(4.2), Inches(4), Inches(0.4), 13, color=GRAY, align=PP_ALIGN.CENTER)

    txt(s, "11% → 87.5%",
        Inches(1), Inches(5.3), Inches(11.3), Inches(0.7), 22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(s, "하이브리드 접근으로 8배 개선 · 진짜 매치를 놓친 경우(FN) 0건",
        Inches(1), Inches(6.0), Inches(11.3), Inches(0.5), 14, color=MED, align=PP_ALIGN.CENTER, italic=True)

    # ── Slide 9: Insights ───────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    txt(s, "08 · Insights",
        Inches(1), Inches(0.8), Inches(11), Inches(0.4), 12, color=ACCENT, bold=True)
    txt(s, "바이브코딩이 남긴 것",
        Inches(1), Inches(1.3), Inches(11), Inches(1.0), 32, bold=True, color=DARK)

    insights = [
        ("🧠  Claude 활용 전략",
         "설계는 Claude.ai에서 천천히 · 코드 실행은 Claude Code로 빠르게 · 역할 분리"),
        ("📐  설계 > 코드",
         "4단계 파이프라인 구조를 먼저 확정 → 삽질 80% 줄어듦"),
        ("✅  정답 데이터 라벨링의 힘",
         "AI 결과 63명을 내가 직접 라벨링 → 어디서 틀렸는지 보여야 프롬프트 개선 가능"),
        ("🔑  \"좋은 모델\"보다 \"좋은 입력\"",
         "같은 AI라도 프롬프트 텍스트 하나 추가로 정확도 11% → 87.5%"),
        ("🌍  기술보다 환경이 문제",
         "NAMC는 사내망이라 클라우드 배포 자체가 불가능 — 일찍 알았다면 처음부터 로컬로"),
    ]
    for i, (title, desc) in enumerate(insights):
        top = Inches(2.5 + i * 0.85)
        txt(s, title, Inches(1), top, Inches(11), Inches(0.4), 15, bold=True, color=DARK)
        txt(s, desc,  Inches(1.4), top + Inches(0.4), Inches(11), Inches(0.4), 12, color=MED)

    # ── Slide 10: Next Steps ────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    txt(s, "09 · Next Steps",
        Inches(1), Inches(0.8), Inches(11), Inches(0.4), 12, color=ACCENT, bold=True)
    txt(s, "앞으로",
        Inches(1), Inches(1.3), Inches(11), Inches(1.0), 32, bold=True, color=DARK)

    bullets(s, [
        "스토리 이미지 판별 정확도 개선 (24h 만료 케이스 대응)",
        "주간 · 월간 자동 스케줄링 (Cron + Slack 알림)",
        "팀원용 로컬 배포 가이드 및 도구 표준화",
        "NAMC 외부 접근 가능성 사내 협의 (장기 과제)",
        "CSV · Google Sheets 자동 리포트 (현재 완료)",
    ], Inches(1), Inches(2.6), Inches(11.3), Inches(4), size=17)

    txt(s, "감사합니다 🙌",
        Inches(1), Inches(6.4), Inches(11.3), Inches(0.5), 18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Save
    out = "presentation_vibecoding_20260418.pptx"
    prs.save(out)
    import os
    print(f"✅ 저장: {out} ({os.path.getsize(out)//1024} KB · {len(prs.slides)} 슬라이드)")


if __name__ == "__main__":
    make()
